#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT页码切分工具
用于从指定PPT文件中提取特定页码范围的幻灯片并保存为新的PPT文件
支持同时处理同名的xlsx文件
"""

import argparse
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import json
import shutil
def find_and_process_xlsx(ppt_path, start_page, end_page, output_dir):
    """
    在PPT文件同一路径下查找同名的xlsx文件并输出内容到终端
    
    Args:
        ppt_path (str): PPT文件路径
        start_page (int): 起始页码
        end_page (int): 结束页码
        output_dir (Path): 输出目录（此函数中不使用）
    
    Returns:
        bool: 是否找到并处理了xlsx文件
    """
    
    # 构造同名xlsx文件路径
    ppt_file = Path(ppt_path)
    xlsx_path = ppt_file.parent / f"{ppt_file.stem}.xlsx"
    
    if not xlsx_path.exists():
        print(f"未找到同名xlsx文件: {xlsx_path}")
        return
    
    # 读取xlsx文件的第一个工作表
    df = pd.read_excel(xlsx_path, sheet_name=0)  # 读取第一个工作表
        
    # 输出内容到终端
    print(f"📊 xlsx文件第一个工作表 ({len(df)} 行 x {len(df.columns)} 列)")
    print("-" * 40)
    
    
    # 过滤数据：根据第二列的序号筛选
    filtered_rows = []
    
    for idx, row in df.iterrows():
        row_data = [str(val) for val in row.values]
        
        # 检查第二列是否存在且为有效的序号
        if len(row_data) >= 2:
            try:
                # 尝试将第二列转换为整数
                page_num = int(float(row_data[1]))  # 先转float再转int，处理可能的小数点
                
                # 检查是否在指定范围内
                if start_page <= page_num <= end_page:
                    filtered_rows.append({
                        "id": row_data[0],
                        "page": page_num,
                        "content": row_data[2]
                    })
            except (ValueError, TypeError):
                # 如果第二列不是数字，跳过这一行
                continue
    
    return filtered_rows

def extract_slides(input_path, start_page, end_page, output_path=None):
    """
    从PPT文件中提取指定页码范围的幻灯片
    
    Args:
        input_path (str): 输入PPT文件路径
        start_page (int): 起始页码（从1开始）
        end_page (int): 结束页码（从1开始，包含）
        output_path (str, optional): 输出文件路径，默认为None时自动生成到./buffer目录
    
    Returns:
        dict: 包含输出文件路径的字典 {'ppt': ppt_path, 'xlsx_found': bool}
    """
    
    # 验证输入文件是否存在
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"输入文件不存在: {input_path}")
    
    # 生成输出文件路径
    if output_path is None:
        # 确保buffer目录存在
        buffer_dir = Path("./buffer")
        buffer_dir.mkdir(exist_ok=True)
        
        input_file = Path(input_path)
        output_path = buffer_dir / f"{input_file.stem}_{start_page}-{end_page}{input_file.suffix}"
    else:
        buffer_dir = Path(output_path).parent
        buffer_dir.mkdir(exist_ok=True)
    
    shutil.copy2(input_path, output_path)
    
    # 加载复制的PPT文件
    try:
        prs = Presentation(output_path)
    except Exception as e:
        raise ValueError(f"无法读取PPT文件: {e}")
    
    total_slides = len(prs.slides)
    
    # 验证页码范围
    if start_page < 1 or end_page < 1:
        raise ValueError("页码必须从1开始")
    
    if start_page > total_slides or end_page > total_slides:
        raise ValueError(f"页码超出范围，PPT总共有{total_slides}页")
    
    if start_page > end_page:
        raise ValueError("起始页码不能大于结束页码")
    
    slides_to_remove = []
    
    for i in range(total_slides):
        slide_index = i + 1  # 转换为1-based索引
        if slide_index < start_page or slide_index > end_page:
            slides_to_remove.append(i)
    
    for slide_index in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[slide_index].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_index]
    
    # 保存修改后的PPT
    try:
        prs.save(output_path)
        ppt_output = str(output_path)
        print(f"💾 已保存提取的幻灯片到: {ppt_output}")
    except Exception as e:
        raise ValueError(f"保存文件失败: {e}")
    
    xlsx_found = find_and_process_xlsx(input_path, start_page, end_page, buffer_dir)

    result_data = {
        'ppt': ppt_output,
        'scripts': xlsx_found
    }

    json_path = Path(ppt_output).with_suffix('.json')
    try:
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(result_data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"保存JSON文件失败: {e}")

    return result_data

def main():

    parser = argparse.ArgumentParser()
    
    parser.add_argument('input_path', help='输入PPT文件路径')    
    parser.add_argument('start_page', type=int, help='起始页码')
    parser.add_argument('end_page', type=int, help='结束页码')
    parser.add_argument('-o', '--output', dest='output_path', help='输出PPT文件路径（可选，默认自动生成）')
    
    args = parser.parse_args()
    
    start_page = args.start_page
    end_page = args.end_page
    
    result = extract_slides(
        args.input_path, 
        start_page, 
        end_page, 
        args.output_path
    )
    
    # 显示PPT处理结果
    ppt_output = result['ppt']
    print(f"📁 PPT输出文件: {ppt_output}")
    
    # 显示xlsx处理结果
    xlsx_found = result['scripts']
    print(xlsx_found)

    print("🎉 处理完成！")
        

if __name__ == "__main__":
    main()
