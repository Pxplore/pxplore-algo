from pptx import Presentation
import re
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import jieba
import argparse
import sys
import os

# TODO: 简单的测试脚本，需要进一步优化

# 关键词过滤
COVER_KEYWORDS = ["welcome", "课程", "讲座", "title", "封面", "欢迎", "介绍"]
TOC_KEYWORDS = ["目录", "contents", "提纲", "大纲", "outline", "agenda"]
TRANSITION_KEYWORDS = ["小结", "回顾", "总结", "结论", "思考", "练习", "作业", "thank", "谢谢", "q&a", "问答"]
SECTION_KEYWORDS = ["第一章", "第二章", "第三章", "第四章", "第五章", "chapter", "section", "part", "模块", "单元"]

def extract_all_text(slide):
    """提取页面中的所有文本内容"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
    return texts

def extract_slide_content(slide):
    """从一页中提取详细内容信息"""
    all_texts = extract_all_text(slide)
    if not all_texts:
        return {"title": "", "content": "", "all_text": "", "text_count": 0}
    
    # 假设第一个非空文本是标题，其余是内容
    title = all_texts[0] if all_texts else ""
    content = " ".join(all_texts[1:]) if len(all_texts) > 1 else ""
    all_text = " ".join(all_texts)
    
    return {
        "title": title,
        "content": content,
        "all_text": all_text,
        "text_count": len(all_text),
        "text_list": all_texts
    }

def extract_title_number(title):
    """提取标题中的序号信息"""
    if not title:
        return None, None
    
    title = title.strip()
    
    # 匹配各种编号格式
    patterns = [
        r'^(\d+)\.(\d+)\.(\d+)',  # 1.2.3
        r'^(\d+)\.(\d+)',         # 1.2
        r'^(\d+)\.',              # 1.
        r'^(\d+)\s',              # 1 空格
        r'第(\d+)章',              # 第1章
        r'第(\d+)节',              # 第1节
        r'第([一二三四五六七八九十]+)章',  # 第一章
        r'Chapter\s*(\d+)',       # Chapter 1
        r'Section\s*(\d+)',       # Section 1
        r'^([一二三四五六七八九十]+)、',  # 一、
        r'^\((\d+)\)',            # (1)
        r'^([A-Z])\.',            # A.
        r'^([a-z])\.',            # a.
    ]
    
    for pattern in patterns:
        match = re.search(pattern, title, re.IGNORECASE)
        if match:
            numbers = match.groups()
            # 转换中文数字
            chinese_nums = {'一': 1, '二': 2, '三': 3, '四': 4, '五': 5, 
                          '六': 6, '七': 7, '八': 8, '九': 9, '十': 10}
            converted_nums = []
            for num in numbers:
                if num in chinese_nums:
                    converted_nums.append(chinese_nums[num])
                elif num.isdigit():
                    converted_nums.append(int(num))
                elif num.isalpha():
                    converted_nums.append(ord(num.upper()) - ord('A') + 1)
            return tuple(converted_nums), pattern
    
    return None, None

def detect_title_level(title):
    """检测标题层次级别"""
    if not title:
        return 0
    
    title_lower = title.lower()
    numbers, pattern = extract_title_number(title)
    
    # 基于编号判断层次
    if numbers:
        if len(numbers) >= 3:  # 1.2.3
            return 3
        elif len(numbers) == 2:  # 1.2
            return 2
        elif len(numbers) == 1:
            if pattern and ('章' in pattern or 'chapter' in pattern.lower()):
                return 1
            else:
                return 2
    
    # 检测章节标题
    if any(keyword in title_lower for keyword in SECTION_KEYWORDS):
        return 1
    
    # 根据文本长度判断
    if len(title) > 30:
        return 3  # 可能是内容而非标题
    elif len(title) > 15:
        return 2
    else:
        return 1

def is_cover_page(slide_content, page_index):
    """改进的封面页识别"""
    if page_index == 0:
        return True
    
    title = slide_content["title"].lower()
    all_text = slide_content["all_text"].lower()
    
    # 检查关键词
    if any(k in all_text for k in COVER_KEYWORDS):
        return True
    
    # 检查文本量（封面页通常文字较少）
    if slide_content["text_count"] < 50 and any(k in title for k in ["welcome", "欢迎", "介绍"]):
        return True
    
    return False

def is_toc_page(slide_content):
    """改进的目录页识别"""
    title = slide_content["title"].lower()
    all_text = slide_content["all_text"].lower()
    
    # 检查关键词
    if any(k in title or k in all_text for k in TOC_KEYWORDS):
        return True
    
    # 检查是否有大量编号列表（目录特征）
    text_lines = slide_content["text_list"]
    numbered_lines = sum(1 for line in text_lines if re.match(r'^\d+\.|\d+\s', line.strip()))
    if len(text_lines) > 3 and numbered_lines / len(text_lines) > 0.5:
        return True
    
    return False

def is_transition_page(slide_content):
    """改进的过渡页识别"""
    title = slide_content["title"].lower()
    all_text = slide_content["all_text"].lower()
    
    # 检查关键词
    if any(k in title or k in all_text for k in TRANSITION_KEYWORDS):
        return True
    
    # 检查文本量（过渡页通常文字很少）
    if slide_content["text_count"] < 30:
        return True
    
    return False

def calculate_content_similarity(content1, content2):
    """计算两个页面内容的相似度"""
    if not content1 or not content2:
        return 0
    
    # 中文分词
    def tokenize(text):
        return " ".join(jieba.cut(text))
    
    texts = [tokenize(content1), tokenize(content2)]
    
    try:
        vectorizer = TfidfVectorizer(max_features=1000, stop_words=None)
        tfidf_matrix = vectorizer.fit_transform(texts)
        similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
        return similarity
    except:
        return 0

def clean_title(text):
    """清洗标题文本，去掉换行等符号"""
    if not text:
        return ""
    return re.sub(r'\s+', ' ', text.strip())

def should_split_by_numbering(current_slide, prev_slide):
    """基于编号判断是否应该切分"""
    curr_nums, curr_pattern = extract_title_number(current_slide["title"])
    prev_nums, prev_pattern = extract_title_number(prev_slide["title"])
    
    if not curr_nums or not prev_nums:
        return False
    
    # 只有主编号变化才强制切分（更宽松的条件）
    if len(curr_nums) >= 1 and len(prev_nums) >= 1:
        if curr_nums[0] != prev_nums[0]:  # 主编号变化
            return True
    
    return False

def smart_segment_slides(content_slides, similarity_threshold=0.2):
    """智能切分页面，返回连续区间，更注重连贯性"""
    if not content_slides:
        return []
    
    if len(content_slides) == 1:
        return [[content_slides[0]["index"], content_slides[0]["index"]]]
    
    split_points = [0]  # 切分点索引
    
    for i in range(1, len(content_slides)):
        current_slide = content_slides[i]
        prev_slide = content_slides[i-1]
        
        should_split = False
        
        # 1. 基于编号的强制切分（只有主编号变化）
        if should_split_by_numbering(current_slide, prev_slide):
            should_split = True
        
        # 2. 章节级别标题变化（更严格的条件）
        elif current_slide["title"] != prev_slide["title"] and current_slide["title"]:
            current_level = detect_title_level(current_slide["title"])
            prev_level = detect_title_level(prev_slide["title"])
            
            # 只有明显的上级标题变化才切分
            if current_level == 1 or (current_level <= prev_level and current_level <= 2):
                # 检查是否真的是重要的标题变化
                curr_nums, _ = extract_title_number(current_slide["title"])
                if curr_nums or current_level == 1:  # 有编号或是章节级标题
                    should_split = True
        
        # 3. 页面间隔过大（减少阈值）
        if not should_split:
            page_gap = current_slide["index"] - prev_slide["index"]
            if page_gap > 5:  # 增加阈值，减少碎片化
                should_split = True
        
        # 4. 内容相似度过低（提高阈值，减少切分）
        if not should_split:
            similarity = calculate_content_similarity(
                prev_slide["all_text"], 
                current_slide["all_text"]
            )
            if similarity < similarity_threshold and len(current_slide["all_text"]) > 100:
                # 只有在内容足够多且相似度很低时才切分
                should_split = True
        
        if should_split:
            split_points.append(i)
    
    # 转换为连续区间
    snippets = []
    for i in range(len(split_points)):
        start_idx = split_points[i]
        end_idx = split_points[i + 1] - 1 if i + 1 < len(split_points) else len(content_slides) - 1
        
        start_page = content_slides[start_idx]["index"]
        end_page = content_slides[end_idx]["index"]
        
        snippets.append([start_page, end_page])
    
    # 后处理：合并过小的片段
    final_snippets = []
    i = 0
    while i < len(snippets):
        current_snippet = snippets[i]
        snippet_length = current_snippet[1] - current_snippet[0] + 1
        
        # 如果当前片段太小（少于3页），尝试合并
        if snippet_length < 3 and final_snippets:
            # 合并到前一个片段
            final_snippets[-1][1] = current_snippet[1]
        elif snippet_length < 3 and i + 1 < len(snippets):
            # 合并到下一个片段
            snippets[i + 1][0] = current_snippet[0]
        else:
            final_snippets.append(current_snippet)
        i += 1
    
    return final_snippets

def extract_ppt_structure(ppt_path, similarity_threshold=0.2):
    """主函数：提取PPT结构并进行智能切分"""
    prs = Presentation(ppt_path)
    slides = prs.slides

    slide_info = []
    toc_page_index = None
    
    for idx, slide in enumerate(slides):
        content = extract_slide_content(slide)
        content["title"] = clean_title(content["title"])
        
        info = {
            "index": idx,
            "title": content["title"],
            "content": content["content"],
            "all_text": content["all_text"],
            "text_count": content["text_count"],
            "title_level": detect_title_level(content["title"]),
            "title_number": extract_title_number(content["title"])[0],
            "is_cover": is_cover_page(content, idx),
            "is_toc": is_toc_page(content),
            "is_transition": is_transition_page(content)
        }
        slide_info.append(info)
        
        # 记录目录页位置
        if info["is_toc"] and toc_page_index is None:
            toc_page_index = idx
    
    # 过滤规则：
    # 1. 目录页之前的所有页面都不保留
    # 2. 过滤封面页、目录页、过渡页
    start_index = toc_page_index + 1 if toc_page_index is not None else 0
    content_slides = []
    
    for info in slide_info:
        if (info["index"] >= start_index and 
            not (info["is_cover"] or info["is_toc"] or info["is_transition"])):
            content_slides.append(info)
    
    # 智能切分
    snippets = smart_segment_slides(content_slides, similarity_threshold)
    
    return snippets, slide_info

# 示例用法
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='PPT智能切分工具')
    parser.add_argument('ppt_path', help='PPT文件路径')
    parser.add_argument('--similarity', '-s', type=float, default=0.2, 
                       help='内容相似度阈值 (0-1, 默认0.2, 值越小越倾向于切分)')
    parser.add_argument('--output', '-o', help='输出结果到文件')
    parser.add_argument('--verbose', '-v', action='store_true', help='显示详细信息')
    
    args = parser.parse_args()
    
    # 检查文件是否存在
    if not os.path.exists(args.ppt_path):
        print(f"错误: 文件 '{args.ppt_path}' 不存在")
        sys.exit(1)
    
    # 检查文件扩展名
    if not args.ppt_path.lower().endswith(('.ppt', '.pptx')):
        print(f"警告: 文件 '{args.ppt_path}' 可能不是PPT文件")
    
    try:
        print(f"正在处理文件: {args.ppt_path}")
        print(f"相似度阈值: {args.similarity}")
        print("-" * 50)
        
        snippets, all_info = extract_ppt_structure(args.ppt_path, args.similarity)
        
        # 输出结果
        result_lines = []
        result_lines.append("智能切片结果：")
        for i, s in enumerate(snippets):
            line = f"Snippet {i+1}: [{s[0]}, {s[1]}]"
            result_lines.append(line)
            print(line)
        
        if args.verbose:
            print("\n页面详情：")
            result_lines.append("\n页面详情：")
            for info in all_info:
                numbers = info.get('title_number', None)
                number_str = f"编号:{numbers}" if numbers else "无编号"
                detail = (f"页 {info['index']}: {info['title'][:30]}... "
                         f"(级别:{info['title_level']}, {number_str}, 文字数:{info['text_count']}, "
                         f"封面:{info['is_cover']}, 目录:{info['is_toc']}, 过渡:{info['is_transition']})")
                result_lines.append(detail)
                print(detail)
        
        # 输出到文件
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                f.write('\n'.join(result_lines))
            print(f"\n结果已保存到: {args.output}")
        
        print(f"\n总共生成 {len(snippets)} 个切片")
        
    except Exception as e:
        print(f"处理文件时出错: {str(e)}")
        sys.exit(1)